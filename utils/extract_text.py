from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
from PIL import Image
from bs4 import BeautifulSoup
import markdown
import fitz
import pytesseract
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
def extract_text(file):

    text = ""

    if file.name.endswith(".txt"):
        text = file.read().decode("utf-8")

    elif file.name.endswith(".pdf"):
        file.seek(0)
        pdf = fitz.open(
            stream=file.read(),
            filetype="pdf"
        )
        for page in pdf:
            text += page.get_text() + "\n"
    elif file.name.endswith(".docx"):
        doc = Document(file)

        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file.name.endswith(".pptx"):

        ppt = Presentation(file)

        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file.name.endswith(".csv"):
        df = pd.read_csv(file)
        text = df.to_string()
    elif file.name.endswith(".xlsx"):
        df = pd.read_excel(file)
        text = df.to_string(index=False)
    elif file.name.endswith((".jpg", ".jpeg", ".png")):
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
    elif file.name.endswith(".html"):
        html_content = file.read().decode("utf-8")
        soup = BeautifulSoup(html_content, "html.parser")
        text = soup.get_text(separator="\n")
    elif file.name.endswith(".md"):
        md_content = file.read().decode("utf-8")
        html = markdown.markdown(md_content)
        soup = BeautifulSoup(html, "html.parser")
        text = soup.get_text(separator="\n")
    return text